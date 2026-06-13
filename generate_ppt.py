import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------
# CONSTANTS & COLOR PALETTE (Modern Slate/Indigo Theme)
# ---------------------------------------------------------
NAVY = RGBColor(11, 19, 43)        # Primary dark theme background
SLATE = RGBColor(15, 23, 42)       # Text dark or cards dark
INDIGO = RGBColor(48, 127, 226)    # Primary accent color
GREEN = RGBColor(22, 163, 74)      # Secondary accent/success color
GOLD = RGBColor(250, 204, 21)      # Gold highlights
LIGHT_BG = RGBColor(241, 245, 249) # Very light gray for content backgrounds
CARD_BG = RGBColor(255, 255, 255)  # Pure white for content cards
DARK_TEXT = RGBColor(30, 41, 59)   # Slate-800 for primary text
LIGHT_TEXT = RGBColor(241, 245, 249)# Off-white for dark background text
MUTED_TEXT = RGBColor(100, 116, 139)# Slate-500 for captions / secondary text

FONT_TITLE = "Outfit"
FONT_BODY = "Inter"

# ---------------------------------------------------------
# HELPER FUNCTIONS FOR CLEAN LAYOUTS
# ---------------------------------------------------------
def create_blank_slide(prs):
    """Creates a completely blank slide layout."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def add_solid_background(slide, color):
    """Draws a full-slide rectangle to establish background color."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background() # No border
    return bg

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=None):
    """Draws a card background block to group content sections."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card

def add_text_box(slide, left, top, width, height, text, size=14, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font_name=FONT_BODY):
    """Helper to add clean, styled text boxes quickly."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_point(tf, bold_prefix, text, size=14, color=DARK_TEXT):
    """Adds a bullet point with a bold prefix to a text frame."""
    p = tf.add_paragraph()
    p.space_after = Pt(8)
    p.font.size = Pt(size)
    p.font.name = FONT_BODY
    
    run1 = p.add_run()
    run1.text = "• " + bold_prefix + ": "
    run1.font.bold = True
    run1.font.color.rgb = color
    
    run2 = p.add_run()
    run2.text = text
    run2.font.bold = False
    run2.font.color.rgb = color

def create_slide_header(slide, title, category="INTERNSHIP PROJECT"):
    """Adds a standard premium header to content slides."""
    # Category / Tag
    add_text_box(slide, 0.8, 0.4, 11.7, 0.3, category.upper(), size=10, color=INDIGO, bold=True, font_name=FONT_TITLE)
    # Slide Title
    add_text_box(slide, 0.8, 0.7, 11.7, 0.6, title, size=28, color=NAVY, bold=True, font_name=FONT_TITLE)
    # Accent Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = INDIGO
    line.line.fill.background()

# ---------------------------------------------------------
# PRESENTATION GENERATION
# ---------------------------------------------------------
def generate_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # ---------------------------------------------------------
    s1 = create_blank_slide(prs)
    add_solid_background(s1, NAVY)

    # Decorative shape / accent on the left
    decor = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), Inches(7.5))
    decor.fill.solid()
    decor.fill.fore_color.rgb = INDIGO
    decor.line.fill.background()

    # Big Title
    add_text_box(s1, 1.2, 1.8, 11.0, 1.0, "PEOPLEFLOW", size=54, color=GOLD, bold=True, font_name=FONT_TITLE)
    add_text_box(s1, 1.2, 2.7, 11.0, 0.8, "A Modern Enterprise HRMS & Leave Management Platform", size=22, color=LIGHT_TEXT, bold=False, font_name=FONT_TITLE)

    # Accent line
    line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.6), Inches(4.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = INDIGO
    line.line.fill.background()

    # Details Grid
    # Col 1: Student & College Details
    c1_box = s1.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(5.5), Inches(2.5))
    tf1 = c1_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
    
    p1 = tf1.paragraphs[0]
    p1.text = "DEVELOPED BY:"
    p1.font.bold = True
    p1.font.size = Pt(11)
    p1.font.color.rgb = MUTED_TEXT
    p1.font.name = FONT_TITLE
    p1.space_after = Pt(4)
    
    p2 = tf1.add_paragraph()
    p2.text = "Krishna Kant Tiwari"
    p2.font.bold = True
    p2.font.size = Pt(18)
    p2.font.color.rgb = LIGHT_TEXT
    p2.font.name = FONT_BODY
    
    p3 = tf1.add_paragraph()
    p3.text = "Course: B.Tech CSIT (1st Year, Sem 2)\nCollege: Symbiosis University of Applied Sciences, Indore"
    p3.font.size = Pt(13)
    p3.font.color.rgb = LIGHT_TEXT
    p3.font.name = FONT_BODY
    p3.space_before = Pt(4)

    # Col 2: Internship & Mentor Details
    c2_box = s1.shapes.add_textbox(Inches(7.2), Inches(4.2), Inches(5.0), Inches(2.5))
    tf2 = c2_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    
    p4 = tf2.paragraphs[0]
    p4.text = "INTERNSHIP HOST & MENTOR:"
    p4.font.bold = True
    p4.font.size = Pt(11)
    p4.font.color.rgb = MUTED_TEXT
    p4.font.name = FONT_TITLE
    p4.space_after = Pt(4)
    
    p5 = tf2.add_paragraph()
    p5.text = "I-Soft Zone Pvt. Ltd."
    p5.font.bold = True
    p5.font.size = Pt(18)
    p5.font.color.rgb = LIGHT_TEXT
    p5.font.name = FONT_BODY
    
    p6 = tf2.add_paragraph()
    p6.text = "Project Mentor: Pranay Gupta\nRole: Director / Mentor"
    p6.font.size = Pt(13)
    p6.font.color.rgb = LIGHT_TEXT
    p6.font.name = FONT_BODY
    p6.space_before = Pt(4)


    # ---------------------------------------------------------
    # SLIDE 2: Project Overview & Motivation (Light Theme)
    # ---------------------------------------------------------
    s2 = create_blank_slide(prs)
    add_solid_background(s2, LIGHT_BG)
    create_slide_header(s2, "Project Overview & Motivation")

    # Column 1: Problem Statement (Card Layout)
    add_card(s2, 0.8, 1.7, 5.6, 5.0)
    add_text_box(s2, 1.2, 2.0, 4.8, 0.5, "THE PROBLEM", size=18, color=RGBColor(220, 38, 38), bold=True, font_name=FONT_TITLE)
    
    p_box1 = s2.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(4.8), Inches(3.8))
    tf_p1 = p_box1.text_frame
    tf_p1.word_wrap = True
    tf_p1.margin_left = tf_p1.margin_top = tf_p1.margin_right = tf_p1.margin_bottom = 0
    
    add_bullet_point(tf_p1, "Administrative Friction", "Traditional systems require heavy manual overhead for payroll, attendance, and leave management.", color=DARK_TEXT)
    add_bullet_point(tf_p1, "Scattered Records", "Employee profiles, department directories, and company documents exist in silos (spreadsheets, emails).", color=DARK_TEXT)
    add_bullet_point(tf_p1, "Unmanaged Assets", "No structured logs for IT asset assignments, causing lost inventory and delayed handovers.", color=DARK_TEXT)
    add_bullet_point(tf_p1, "Lack of Transparency", "Approvals are delayed due to offline request pipelines, hurting employee satisfaction.", color=DARK_TEXT)

    # Column 2: The Solution (Card Layout)
    add_card(s2, 6.9, 1.7, 5.6, 5.0)
    add_text_box(s2, 7.3, 2.0, 4.8, 0.5, "THE PEOPLEFLOW SOLUTION", size=18, color=GREEN, bold=True, font_name=FONT_TITLE)

    p_box2 = s2.shapes.add_textbox(Inches(7.3), Inches(2.6), Inches(4.8), Inches(3.8))
    tf_p2 = p_box2.text_frame
    tf_p2.word_wrap = True
    tf_p2.margin_left = tf_p2.margin_top = tf_p2.margin_right = tf_p2.margin_bottom = 0
    
    add_bullet_point(tf_p2, "Unified Dashboard", "All operations (Attendance, Leaves, Assets, Analytics) consolidated under one portal.", color=DARK_TEXT)
    add_bullet_point(tf_p2, "Automated Tracking", "One-click digital check-in/check-out with real-time logging and analytics.", color=DARK_TEXT)
    add_bullet_point(tf_p2, "Streamlined Workflows", "Multi-stage digital leave requests and approval pipelines for managers and HR.", color=DARK_TEXT)
    add_bullet_point(tf_p2, "Audited Environment", "Full audit logging for actions (logins, updates, role changes) to ensure security.", color=DARK_TEXT)


    # ---------------------------------------------------------
    # SLIDE 3: System Architecture & Tech Stack (Light Theme)
    # ---------------------------------------------------------
    s3 = create_blank_slide(prs)
    add_solid_background(s3, LIGHT_BG)
    create_slide_header(s3, "System Architecture & Tech Stack")

    # 3 columns for Frontend, Backend, Database
    card_w = 3.6
    card_h = 5.0
    start_x = 0.8
    spacing = 0.45

    # Card 1: Frontend
    x1 = start_x
    add_card(s3, x1, 1.7, card_w, card_h)
    add_text_box(s3, x1 + 0.3, 2.0, card_w - 0.6, 0.4, "FRONTEND CLIENT", size=18, color=INDIGO, bold=True, font_name=FONT_TITLE)
    
    tf_f1 = s3.shapes.add_textbox(Inches(x1 + 0.3), Inches(2.6), Inches(card_w - 0.6), Inches(3.8)).text_frame
    tf_f1.word_wrap = True
    tf_f1.margin_left = tf_f1.margin_top = tf_f1.margin_right = tf_f1.margin_bottom = 0
    add_bullet_point(tf_f1, "React.js & TS", "Typed components for highly scalable code structure.", color=DARK_TEXT)
    add_bullet_point(tf_f1, "Tailwind CSS", "Fully custom responsive design utilities.", color=DARK_TEXT)
    add_bullet_point(tf_f1, "Vite", "Next-gen lightning-fast frontend compilation.", color=DARK_TEXT)
    add_bullet_point(tf_f1, "Framer Motion", "Micro-interactions and fluid animations.", color=DARK_TEXT)

    # Card 2: Backend
    x2 = x1 + card_w + spacing
    add_card(s3, x2, 1.7, card_w, card_h)
    add_text_box(s3, x2 + 0.3, 2.0, card_w - 0.6, 0.4, "BACKEND SERVICES", size=18, color=INDIGO, bold=True, font_name=FONT_TITLE)

    tf_f2 = s3.shapes.add_textbox(Inches(x2 + 0.3), Inches(2.6), Inches(card_w - 0.6), Inches(3.8)).text_frame
    tf_f2.word_wrap = True
    tf_f2.margin_left = tf_f2.margin_top = tf_f2.margin_right = tf_f2.margin_bottom = 0
    add_bullet_point(tf_f2, "Node.js & Express", "Asynchronous, fast REST API controllers.", color=DARK_TEXT)
    add_bullet_point(tf_f2, "Security & Auth", "JWT-based sessions, Bcrypt hashing, Helmet protection.", color=DARK_TEXT)
    add_bullet_point(tf_f2, "Middleware", "Role validation, rate-limiting, and error-handling filters.", color=DARK_TEXT)
    add_bullet_point(tf_f2, "Nodemailer", "SMTP integration for emails.", color=DARK_TEXT)

    # Card 3: Database & Cloud
    x3 = x2 + card_w + spacing
    add_card(s3, x3, 1.7, card_w, card_h)
    add_text_box(s3, x3 + 0.3, 2.0, card_w - 0.6, 0.4, "DATABASE & CLOUD", size=18, color=INDIGO, bold=True, font_name=FONT_TITLE)

    tf_f3 = s3.shapes.add_textbox(Inches(x3 + 0.3), Inches(2.6), Inches(card_w - 0.6), Inches(3.8)).text_frame
    tf_f3.word_wrap = True
    tf_f3.margin_left = tf_f3.margin_top = tf_f3.margin_right = tf_f3.margin_bottom = 0
    add_bullet_point(tf_f3, "PostgreSQL", "Robust, relational schema holding transactions and profiles.", color=DARK_TEXT)
    add_bullet_point(tf_f3, "Render", "Production host for Node backend server.", color=DARK_TEXT)
    add_bullet_point(tf_f3, "Vercel", "Production host for static built frontend application.", color=DARK_TEXT)
    add_bullet_point(tf_f3, "Automatic Migrations", "Database schemas sync automatically on boot.", color=DARK_TEXT)


    # ---------------------------------------------------------
    # SLIDE 4: Core Product Features (Light Theme)
    # ---------------------------------------------------------
    s4 = create_blank_slide(prs)
    add_solid_background(s4, LIGHT_BG)
    create_slide_header(s4, "Core Product Features")

    # 4 cards in a 2x2 grid
    cw = 5.6
    ch = 2.3
    y_starts = [1.7, 4.3]
    x_starts = [0.8, 6.9]

    # Grid 1: Employee Management
    add_card(s4, x_starts[0], y_starts[0], cw, ch)
    add_text_box(s4, x_starts[0] + 0.3, y_starts[0] + 0.2, cw - 0.6, 0.3, "EMPLOYEE DIRECTORY & PROFILES", size=16, color=INDIGO, bold=True, font_name=FONT_TITLE)
    tf_g1 = s4.shapes.add_textbox(Inches(x_starts[0] + 0.3), Inches(y_starts[0] + 0.6), Inches(cw - 0.6), Inches(1.5)).text_frame
    tf_g1.word_wrap = True
    tf_g1.margin_left = tf_g1.margin_top = tf_g1.margin_right = tf_g1.margin_bottom = 0
    add_bullet_point(tf_g1, "Comprehensive Profiles", "Work history, salaries, document store, managers.", color=DARK_TEXT)
    add_bullet_point(tf_g1, "Dynamic Search & Filters", "Query by name, skills, role, or department.", color=DARK_TEXT)

    # Grid 2: Attendance Tracking
    add_card(s4, x_starts[1], y_starts[0], cw, ch)
    add_text_box(s4, x_starts[1] + 0.3, y_starts[0] + 0.2, cw - 0.6, 0.3, "DAILY ATTENDANCE MANAGEMENT", size=16, color=INDIGO, bold=True, font_name=FONT_TITLE)
    tf_g2 = s4.shapes.add_textbox(Inches(x_starts[1] + 0.3), Inches(y_starts[0] + 0.6), Inches(cw - 0.6), Inches(1.5)).text_frame
    tf_g2.word_wrap = True
    tf_g2.margin_left = tf_g2.margin_top = tf_g2.margin_right = tf_g2.margin_bottom = 0
    add_bullet_point(tf_g2, "Check-in / Check-out", "Single-click digital registry.", color=DARK_TEXT)
    add_bullet_point(tf_g2, "Analytics Dashboards", "Visual metrics of monthly/weekly check-ins.", color=DARK_TEXT)

    # Grid 3: Leave Management
    add_card(s4, x_starts[0], y_starts[1], cw, ch)
    add_text_box(s4, x_starts[0] + 0.3, y_starts[1] + 0.2, cw - 0.6, 0.3, "LEAVE WORKFLOWS & APPROVALS", size=16, color=INDIGO, bold=True, font_name=FONT_TITLE)
    tf_g3 = s4.shapes.add_textbox(Inches(x_starts[0] + 0.3), Inches(y_starts[1] + 0.6), Inches(cw - 0.6), Inches(1.5)).text_frame
    tf_g3.word_wrap = True
    tf_g3.margin_left = tf_g3.margin_top = tf_g3.margin_right = tf_g3.margin_bottom = 0
    add_bullet_point(tf_g3, "Leave Balance Tracking", "Tracks allocated, used, and pending leaves.", color=DARK_TEXT)
    add_bullet_point(tf_g3, "Multi-level Approvals", "Routed dynamically to managers and HR.", color=DARK_TEXT)

    # Grid 4: Audit & Security
    add_card(s4, x_starts[1], y_starts[1], cw, ch)
    add_text_box(s4, x_starts[1] + 0.3, y_starts[1] + 0.2, cw - 0.6, 0.3, "AUDIT ENGINE & ROLE-BASED ACCESS", size=16, color=INDIGO, bold=True, font_name=FONT_TITLE)
    tf_g4 = s4.shapes.add_textbox(Inches(x_starts[1] + 0.3), Inches(y_starts[1] + 0.6), Inches(cw - 0.6), Inches(1.5)).text_frame
    tf_g4.word_wrap = True
    tf_g4.margin_left = tf_g4.margin_top = tf_g4.margin_right = tf_g4.margin_bottom = 0
    add_bullet_point(tf_g4, "Granular Permissions", "Super Admin, HR, Manager, and Employee tiers.", color=DARK_TEXT)
    add_bullet_point(tf_g4, "Action Log Audits", "Secure trails logging logins, actions, and IP addresses.", color=DARK_TEXT)


    # ---------------------------------------------------------
    # SLIDE 5: Database Schema & Relations (Light Theme)
    # ---------------------------------------------------------
    s5 = create_blank_slide(prs)
    add_solid_background(s5, LIGHT_BG)
    create_slide_header(s5, "Database Schema & Relations")

    # Content Box for Database schema overview
    add_card(s5, 0.8, 1.7, 11.7, 5.0)

    # Left: Database Entities Explanation
    tf_db = s5.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(5.0), Inches(4.4)).text_frame
    tf_db.word_wrap = True
    tf_db.margin_left = tf_db.margin_top = tf_db.margin_right = tf_db.margin_bottom = 0
    
    add_text_box(s5, 1.2, 2.0, 5.0, 0.4, "PRIMARY DATA ENTITIES", size=18, color=NAVY, bold=True, font_name=FONT_TITLE)
    
    # Simple list of main tables and purposes
    tf_tbls = s5.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(5.0), Inches(3.8)).text_frame
    tf_tbls.word_wrap = True
    add_bullet_point(tf_tbls, "users", "Credentials, status, locked status, and registration info.", color=DARK_TEXT)
    add_bullet_point(tf_tbls, "employees", "Demographics, contact info, job level, salary, manager_id.", color=DARK_TEXT)
    add_bullet_point(tf_tbls, "departments", "Stores department details and maps the head manager.", color=DARK_TEXT)
    add_bullet_point(tf_tbls, "leave_requests", "Requested ranges, total days, approvals, and status.", color=DARK_TEXT)
    add_bullet_point(tf_tbls, "assets", "Company inventory with serial codes and current assignments.", color=DARK_TEXT)

    # Right: Relational mapping visual description
    add_text_box(s5, 6.8, 2.0, 5.2, 0.4, "RELATIONAL ARCHITECTURE", size=18, color=NAVY, bold=True, font_name=FONT_TITLE)
    
    tf_rel = s5.shapes.add_textbox(Inches(6.8), Inches(2.5), Inches(5.2), Inches(3.8)).text_frame
    tf_rel.word_wrap = True
    tf_rel.margin_left = tf_rel.margin_top = tf_rel.margin_right = tf_rel.margin_bottom = 0
    add_bullet_point(tf_rel, "1-to-1 Mapping", "users ↔ employees linked by user_id to isolate auth details from profile metadata.", color=DARK_TEXT)
    add_bullet_point(tf_rel, "Self-Referencing Manager Link", "employees.manager_id links to employees.id to resolve hierarchical approvals.", color=DARK_TEXT)
    add_bullet_point(tf_rel, "Junction Tables", "user_roles and role_permissions resolve multi-role assignments cleanly.", color=DARK_TEXT)
    add_bullet_point(tf_rel, "Integrity Rules", "ON DELETE CASCADE on tokens/logs prevents database orphan records.", color=DARK_TEXT)


    # ---------------------------------------------------------
    # SLIDE 6: Key Technical Accomplishments & Fixes (Light Theme)
    # ---------------------------------------------------------
    s6 = create_blank_slide(prs)
    add_solid_background(s6, LIGHT_BG)
    create_slide_header(s6, "Key Technical Accomplishments & Fixes")

    # Column 1: Backend fixes (Card Layout)
    add_card(s6, 0.8, 1.7, 5.6, 5.0)
    add_text_box(s6, 1.2, 2.0, 4.8, 0.4, "BACKEND & INFRASTRUCTURE", size=18, color=INDIGO, bold=True, font_name=FONT_TITLE)
    
    tf_fix1 = s6.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(4.8), Inches(3.9)).text_frame
    tf_fix1.word_wrap = True
    tf_fix1.margin_left = tf_fix1.margin_top = tf_fix1.margin_right = tf_fix1.margin_bottom = 0
    add_bullet_point(tf_fix1, "Cloud PostgreSQL Compatibility", "Upgraded pg database pool config to parse full DATABASE_URL connection strings.", color=DARK_TEXT)
    add_bullet_point(tf_fix1, "SSL Configurations", "Integrated SSL (rejectUnauthorized: false) for cloud database providers.", color=DARK_TEXT)
    add_bullet_point(tf_fix1, "Startup Migrations", "Refactored server.js to run migrations programmatically on startup, ensuring auto-creation of tables on Render.", color=DARK_TEXT)
    add_bullet_point(tf_fix1, "Fault-tolerant Emails", "Made email verification optional and wrapped sendMail in try-catch to prevent mail server down-time from blocking registrations.", color=DARK_TEXT)

    # Column 2: Frontend & Deployment fixes (Card Layout)
    add_card(s6, 6.9, 1.7, 5.6, 5.0)
    add_text_box(s6, 7.3, 2.0, 4.8, 0.4, "FRONTEND & PROXY LAYERS", size=18, color=INDIGO, bold=True, font_name=FONT_TITLE)

    tf_fix2 = s6.shapes.add_textbox(Inches(7.3), Inches(2.5), Inches(4.8), Inches(3.9)).text_frame
    tf_fix2.word_wrap = True
    tf_fix2.margin_left = tf_fix2.margin_top = tf_fix2.margin_right = tf_fix2.margin_bottom = 0
    add_bullet_point(tf_fix2, "Vercel Rewrite Proxy", "Created vercel.json rewrites for /api/* and /uploads/* pointing to Render.", color=DARK_TEXT)
    add_bullet_point(tf_fix2, "CORS Fixes", "Resolved CORS blockages transparently via Vercel proxying and unified allowed headers.", color=DARK_TEXT)
    add_bullet_point(tf_fix2, "React Router Fallback", "Configured Vercel's index.html fallback rewrite to prevent 404 errors when reloading sub-routes.", color=DARK_TEXT)
    add_bullet_point(tf_fix2, "Flexible Client Base URL", "Modified Axios api.ts client to check for VITE_API_URL, matching localhost or cloud hosts dynamically.", color=DARK_TEXT)


    # ---------------------------------------------------------
    # SLIDE 7: Learnings & Conclusion (Dark Theme)
    # ---------------------------------------------------------
    s7 = create_blank_slide(prs)
    add_solid_background(s7, NAVY)

    # Decorative shape / accent on the left
    decor7 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), Inches(7.5))
    decor7.fill.solid()
    decor7.fill.fore_color.rgb = GREEN
    decor7.line.fill.background()

    # Title
    add_text_box(s7, 1.2, 1.0, 11.0, 0.6, "Internship Learnings & Conclusion", size=32, color=GOLD, bold=True, font_name=FONT_TITLE)
    
    # Left column: Tech Learnings
    add_text_box(s7, 1.2, 1.8, 5.0, 0.4, "TECHNICAL ACQUISITIONS", size=18, color=LIGHT_TEXT, bold=True, font_name=FONT_TITLE)
    tf_ln1 = s7.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(5.0), Inches(4.2)).text_frame
    tf_ln1.word_wrap = True
    tf_ln1.margin_left = tf_ln1.margin_top = tf_ln1.margin_right = tf_ln1.margin_bottom = 0
    add_bullet_point(tf_ln1, "State Management", "Hands-on experience in React context and authentication stores.", color=LIGHT_TEXT)
    add_bullet_point(tf_ln1, "Database Transactions", "Mastered SQL scripts, transaction commits, rollbacks, and indexes.", color=LIGHT_TEXT)
    add_bullet_point(tf_ln1, "System Deployments", "Configured production builds, SSL databases, and API reverse proxying.", color=LIGHT_TEXT)
    add_bullet_point(tf_ln1, "Auditing & RBAC", "Designed structured role authentication and security event auditing.", color=LIGHT_TEXT)

    # Right column: Conclusion & Roadmap
    add_text_box(s7, 6.8, 1.8, 5.2, 0.4, "CONCLUSION & FUTURE ROADMAP", size=18, color=LIGHT_TEXT, bold=True, font_name=FONT_TITLE)
    tf_ln2 = s7.shapes.add_textbox(Inches(6.8), Inches(2.3), Inches(5.2), Inches(4.2)).text_frame
    tf_ln2.word_wrap = True
    tf_ln2.margin_left = tf_ln2.margin_top = tf_ln2.margin_right = tf_ln2.margin_bottom = 0
    add_bullet_point(tf_ln2, "Industry Alignment", "Working on PeopleFlow during this internship at I-Soft Zone under Pranay Gupta provided exposure to enterprise coding standards.", color=LIGHT_TEXT)
    add_bullet_point(tf_ln2, "Payroll Integration", "Plan to implement salary structures, tax deductions, and download pay slips.", color=LIGHT_TEXT)
    add_bullet_point(tf_ln2, "AI Shift Planner", "Integrate models to forecast employee schedules and suggest optimization.", color=LIGHT_TEXT)
    add_bullet_point(tf_ln2, "Thank You", "Deep gratitude to SUAS Indore and I-Soft Zone for this learning opportunity.", color=LIGHT_TEXT)

    # Save presentation
    output_filename = "peopleflow_internship_presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved as: {output_filename}")
    return output_filename

if __name__ == "__main__":
    generate_presentation()
